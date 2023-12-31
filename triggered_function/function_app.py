import azure.functions as func
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
import logging
import json
import os
from pptx import Presentation

from dotenv import load_dotenv

load_dotenv()

app = func.FunctionApp()

credentials = DefaultAzureCredential()

def get_url_parts(url):
    # Ersetze '//' durch eine spezielle Zeichenfolge, die wahrscheinlich nicht im String vorhanden ist
    modified_string = url.replace('//', '__double_slash__')
    
    # Teile den String an den einzelnen Schrägstrichen
    parts = modified_string.split('/')
    
    # Wiederherstelle '//' nach dem ersten Element
    parts[0] = parts[0].replace('__double_slash__', '//')
    
    return parts

def update_result_txt(results_def, results):
   blob_service_client = BlobServiceClient(account_url= results_def[0], credential= credentials)
   container_client = blob_service_client.get_container_client(container=results_def[1])

   # Upload results_xx.txt
   results_filename = results_def[2]
   results_f = open(results_filename, "w")
   for element in results:
       results_f.write(element)
   results_f.close()

   #open and read the file after the overwriting:
   results_f = open(results_filename, "r")
   data = results_f.read()
   container_client.upload_blob(name=results_filename, data=data, overwrite=True)
   results_f.close()
   os.remove(results_filename)

def analyze_powerpoint(powerpoint_def, search_string):
    # Read powerpoint data
    blob_service_client = BlobServiceClient(account_url= powerpoint_def[0], credential= credentials)
    container_client = blob_service_client.get_container_client(container=powerpoint_def[1])

    # download blob data 
    blob_client = container_client.get_blob_client(blob= powerpoint_def[2])
    with open(file=os.path.join(r'.', powerpoint_def[2]), mode="wb") as sample_blob:
        download_stream = blob_client.download_blob()
        sample_blob.write(download_stream.readall())

    # Search for String in Powerpoint 
    presentation = Presentation(powerpoint_def[2])

    header = "Searched for " + search_string + " in " + powerpoint_def[2] + "\n" + "\n"
    results = [header] 

    for i, slide in enumerate(presentation.slides):
    
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text

                if search_string in text:
                    lines = text.split("\n")
                    for line in lines:
                        if search_string in line:
                            text_with_slide = "Slide " + str(i+1) + ":     " + line + "\n"
                            results.append(text_with_slide)                  

    if len(results) == 1:
        results.append("String not found in Powerpoint")
    # Remove powerpoint from blob storage
    os.remove(powerpoint_def[2])
    blob_client.delete_blob()
    return results
    

def handle_message(message):
    decoded = json.loads(message)
    results_url = decoded['results_url']
    powerpoint_url = decoded['powerpoint_url'] 
    search_string = decoded['search-string']
    results_def = get_url_parts(results_url)
    powerpoint_def = get_url_parts(powerpoint_url)
    # Analyze
    results = analyze_powerpoint(powerpoint_def, search_string)
    # Update results
    update_result_txt(results_def, results)


@app.queue_trigger(arg_name="azqueue", queue_name="queue-branding-police-app",
                               connection="QueueConnectionString") 
def func_analyze_powerpoint_trigger(azqueue: func.QueueMessage):
    message = azqueue.get_body().decode('utf-8')
    logging.info('Python Queue trigger processed a message: %s',
                message)
    handle_message(message)
